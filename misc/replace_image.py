from pptx.util import Inches


def py_replace_image(pres, label, new_image, new_height=True):
    """Replace a template image with a new image

    Keyword arguments:
     pres -- pptx Pres object containing the text to be replaced
     label -- The unique label of the object containing the text to be replaced
     new_image -- Path to the new image.
     new_height -- Change the height of the new image to match the old image?
    """
    # Get slide index from label
    slide_idx_list = []
    for slide_idx, slide in enumerate(pres.slides):
        target_shape = [shape for shape in slide.shapes if shape.name == label]
        slide_contains_target_shape = len(target_shape) > 0

        if slide_contains_target_shape:
            slide_idx_list.append(slide_idx)

    if len(slide_idx_list) == 1:
        slide_idx = slide_idx_list[0]
    elif len(slide_idx_list) == 0:
        raise Exception("No object with the specified label was found.")
    elif len(slide_idx_list) > 1:
        raise Exception("More than one object with the specified label was found.")

    # Get shape from label
    target_shape = [
        shape for slide in pres.slides for shape in slide.shapes if shape.name == label
    ]

    if len(target_shape) == 0:
        raise Exception("No object with the specified label was found.")
    elif len(target_shape) > 1:
        raise Exception("More than one shape with that label was found.")

    old_image = target_shape[0]

    slide = pres.slides[slide_idx]

    # Make sure the size and position of the new image are the same as for the
    # old image.
    shape_left = Inches(old_image.left.inches)
    shape_top = Inches(old_image.top.inches)
    shape_width = Inches(old_image.width.inches)

    if new_height:
        shape_height = Inches(old_image.height.inches)
    else:
        shape_height = None

    # Add the new image
    new_shape = slide.shapes.add_picture(
        new_image,
        left=shape_left,
        top=shape_top,
        width=shape_width,
        height=shape_height,
    )

    # Delete the old image
    old_image = old_image._element
    new_pic = new_shape._element
    old_image.addnext(new_pic)
    old_image.getparent().remove(old_image)
