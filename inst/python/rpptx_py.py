from pptx.chart.data import CategoryChartData
from pptx.util import Inches

def get_slide_id(pres, slide_num):
  """Get the slide ID for a particular slide in a presentation.
  
  Keyword Arguments:
    pres -- pptx Pres object containing the target shape
    slide_num -- Slide number
  """
  return pres.slides[slide_num - 1].slide_id


def get_shape_with_label(pres, label, slide_num=None):
    """Select the shape with a particular label from a presentation object.

    Keyword Arguments:
      pres -- pptx Pres object containing the target shape
      slide_num -- (Optional) Slide number containing the target shape (if known)
      label -- The unique label of the target shape
    """
    if slide_num is None:
        target_shape = [
            shape
            for slide in pres.slides
            for shape in slide.shapes
            if shape.name == label
        ]
    else:
        slide_num = int(slide_num)
        slide = pres.slides[slide_num - 1]
        target_shape = [shape for shape in slide.shapes if shape.name == label]

    if len(target_shape) == 0:
        raise Exception("No object with the specified label was found.")
    elif len(target_shape) > 1:
        raise Exception("More than one shape with that label was found.")

    return target_shape[0]


def get_slide_index_with_label(pres, label):
    """Find the (zero-indexed) index of the slide in a Presentation object
        containing an object with the target label.

    Args:
        pres (Presentation): A pptx Presentation object
        label (string): The target label

    Returns:
        numeric: Index of the slide containing an object with the target label
    """
    slide_idx_list = []
    for slide_idx, slide in enumerate(pres.slides):
        target_shape = [shape for shape in slide.shapes if shape.name == label]
        slide_contains_target_shape = len(target_shape) > 0

        if slide_contains_target_shape:
            slide_idx_list.append(slide_idx)

    if len(slide_idx_list) == 1:
        return slide_idx_list[0]
    elif len(slide_idx_list) == 0:
        raise Exception("No object with the specified label was found.")
    elif len(slide_idx_list) > 1:
        raise Exception("More than one object with the specified label was found.")


def py_replace_category_plot(pres, label, categories, series_levels, series_values):
    """Replace the data underlying a category plot with the target label

    Args:
        pres (Presentation): A pptx Presentation object
        label (string): The target label
        categories (list): A list containing the levels of the category factor (generally the x-axis variable)
        series_list (list): A list containing list elements of the form `['SeriesName', (1, 2, 3)]`
    """
    # Find shape from label
    old_plot = get_shape_with_label(pres, label)

    new_data = CategoryChartData()
    new_data.categories = categories

    for idx, series in enumerate(series_values):
        if not isinstance(series, list):
            series = [series]
        if not isinstance(series_levels, list):
            series_levels = [series_levels]
        new_data.add_series(series_levels[idx], tuple(series))

    old_plot.chart.replace_data(new_data)


def py_replace_image(pres, label, new_image, new_height=True):
    """Replace a template image with a new image

    Keyword arguments:
     pres -- pptx Pres object containing the text to be replaced
     label -- The unique label of the object containing the text to be replaced
     new_image -- Path to the new image.
     new_height -- Change the height of the new image to match the old image?
    """
    # Get slide index from label
    slide_idx = get_slide_index_with_label(pres, label)

    # Get shape from label
    old_image = get_shape_with_label(pres, label)

    slide = pres.slides[slide_idx]

    # Make sure the size and position of the new image are the same as for the
    # old image.
    shape_left = Inches(old_image.left.inches)
    shape_top = Inches(old_image.top.inches)
    shape_width = Inches(old_image.width.inches)

    if new_height:
        shape_height = None
    else:
        shape_height = Inches(old_image.height.inches)
        

    # Add the new image
    new_shape = slide.shapes.add_picture(
        new_image,
        left=shape_left,
        top=shape_top,
        width=shape_width,
        height=shape_height,
    )

    # Delete the old image
    old_image = old_image._element
    new_pic = new_shape._element
    old_image.addnext(new_pic)
    old_image.getparent().remove(old_image)


def py_replace_table(pres, label, new_table, new_table_shape):
    """Replace text in a text box but retain formatting

    Keyword arguments:
    pres -- pptx Pres object containing the text to be replaced
    label -- The unique label of the object containing the text to be replaced
    new_table -- Tuple containing the contents of the new table elements, in left-to-right,
      top-to-bottom order
    """
    old_table = get_shape_with_label(pres, label)

    # Check the shape of the old table against the shape of the new table
    if (len(old_table.table.rows) != new_table_shape[0]) | (
        len(old_table.table.columns) != new_table_shape[1]
    ):
        err = ["The number of rows and columns in the new table does not match the old table. new: ", str(new_table_shape[0]), " - ", str(new_table_shape[0]), "old: ", str(len(old_table.table.rows)), " - ", str(len(old_table.table.columns))];
        errconcat = " ".join(err);
        print(errconcat)
        raise ValueError(
          errconcat    
        )

    for table_idx, cell in enumerate(old_table.table.iter_cells()):
        # Replace old text with new text, keeping formatting
        paragraph = cell.text_frame.paragraphs[0]

        for run_idx, run in enumerate(paragraph.runs):
            if run_idx == 0:
                continue
            p.remove(run._r)

        paragraph.runs[0].text = new_table[table_idx]


def py_replace_text(pres, label, new_text):
    """Replace text in a text box but retain formatting

    Keyword arguments:
    pres -- pptx Pres object containing the text to be replaced
    label -- The unique label of the object containing the text to be replaced
    new_text -- The new text that the paragraph should contain
    """
    # Get shape from label
    old_text = get_shape_with_label(pres, label)

    paragraph = old_text.text_frame.paragraphs[0]
    p = paragraph._p  # the lxml element containing the `<a:p>` element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text
